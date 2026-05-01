"""Generate CMPE258 Toxic Comment Classification presentation — spacious layout."""

from pathlib import Path
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Paths ──────────────────────────────────────────────────────────────────
REPO = Path(__file__).parent.parent
FIG  = REPO / "report" / "figures"
NB   = REPO / "notebooks"
OUT  = REPO / "report" / "presentation.pptx"

# ── Palette ────────────────────────────────────────────────────────────────
BG        = RGBColor(0x0D, 0x1B, 0x2A)
TITLE_FG  = RGBColor(0x4F, 0xC3, 0xF7)
BODY_FG   = RGBColor(0xE8, 0xEE, 0xF4)
ACCENT    = RGBColor(0xFF, 0xB3, 0x47)
DIM       = RGBColor(0x8A, 0x9B, 0xAD)
TABLE_HDR = RGBColor(0x1A, 0x36, 0x52)
TABLE_ALT = RGBColor(0x12, 0x28, 0x3F)
TABLE_REG = RGBColor(0x0D, 0x1B, 0x2A)
RULE_CLR  = RGBColor(0x4F, 0xC3, 0xF7)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── Content area constants ─────────────────────────────────────────────────
CX = 0.5        # content left x
CY = 1.3        # content top y (below chrome)
CW = 12.33      # content width
CH = 5.9        # content height
CB = 7.1        # content bottom y


# ═══════════════════════════════════════════════════════════════════════════
# PRIMITIVE HELPERS
# ═══════════════════════════════════════════════════════════════════════════

def rect(slide, l, t, w, h, fill=None, line_color=None, line_pt=1):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid() if fill else s.fill.background()
    if fill:
        s.fill.fore_color.rgb = fill
    if line_color:
        s.line.color.rgb = line_color
        s.line.width = Pt(line_pt)
    else:
        s.line.fill.background()
    return s


def txb(slide, text, l, t, w, h,
        size=18, bold=False, italic=False,
        color=None, align=PP_ALIGN.LEFT, wrap=True):
    box = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = wrap
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color or BODY_FG
    return box


def chrome(slide, title, section=""):
    """Dark top bar + rule + title + optional section tag."""
    rect(slide, 0, 0, 13.33, 7.5, fill=BG)
    rect(slide, 0, 0, 13.33, 1.1, fill=RGBColor(0x0A, 0x14, 0x20))
    rect(slide, 0, 1.1, 13.33, 0.045, fill=RULE_CLR)
    txb(slide, title, 0.4, 0.1, 11.5, 0.9,
        size=30, bold=True, color=TITLE_FG)
    if section:
        txb(slide, section, 10.5, 0.2, 2.6, 0.55,
            size=12, color=DIM, align=PP_ALIGN.RIGHT)


def blist(slide, items, l, t, w, h, size=17, gap_pt=5):
    """Bullet list. items = list of (indent_level, text)."""
    box = slide.shapes.add_textbox(
        Inches(l), Inches(t), Inches(w), Inches(h))
    box.word_wrap = True
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for lvl, text in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(gap_pt)
        bullet = "•" if lvl == 0 else "◦"
        pad = "    " * lvl
        r = p.add_run()
        r.text = f"{pad}{bullet}  {text}"
        r.font.size = Pt(size)
        r.font.color.rgb = BODY_FG


def table(slide, headers, rows, l, t, w,
          col_widths=None, hdr_sz=14, row_sz=13,
          row_h=0.44, hdr_h=0.48):
    """Draw table as stacked rects + text boxes."""
    nc = len(headers)
    if col_widths is None:
        col_widths = [w / nc] * nc

    # header
    x = l
    for hdr, cw in zip(headers, col_widths):
        rect(slide, x, t, cw - 0.02, hdr_h, fill=TABLE_HDR)
        txb(slide, hdr, x + 0.06, t + 0.05, cw - 0.14, hdr_h - 0.08,
            size=hdr_sz, bold=True, color=TITLE_FG, align=PP_ALIGN.CENTER)
        x += cw

    # rows
    for ri, row in enumerate(rows):
        rt = t + hdr_h + ri * row_h
        fill = TABLE_ALT if ri % 2 == 0 else TABLE_REG
        x = l
        for cell, cw in zip(row, col_widths):
            rect(slide, x, rt, cw - 0.02, row_h - 0.02, fill=fill)
            bold = str(cell).startswith("**") and str(cell).endswith("**")
            txt  = str(cell).strip("*")
            col  = ACCENT if bold else BODY_FG
            txb(slide, txt, x + 0.06, rt + 0.05,
                cw - 0.14, row_h - 0.1,
                size=row_sz, bold=bold, color=col,
                align=PP_ALIGN.CENTER)
            x += cw

    total_h = hdr_h + len(rows) * row_h
    return t + total_h   # return bottom y


def _clear_cell_borders(tbl):
    """Remove all internal and external cell borders from a native pptx table."""
    for row in tbl.rows:
        for cell in row.cells:
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            for side in ('lnL', 'lnR', 'lnT', 'lnB'):
                tag = qn(f'a:{side}')
                for el in tcPr.findall(tag):
                    tcPr.remove(el)
                ln = etree.SubElement(tcPr, qn(f'a:{side}'))
                etree.SubElement(ln, qn('a:noFill'))


def native_table(slide, headers, rows, l, t, w,
                 col_widths=None, hdr_sz=14, row_sz=13,
                 row_h=0.44, hdr_h=0.48):
    """Native pptx table — single editable shape in PowerPoint.

    Replaces the rect+txb table() helper for slides where in-pptx editing matters.
    """
    n_cols = len(headers)
    n_rows = 1 + len(rows)
    total_h = hdr_h + len(rows) * row_h

    if col_widths is None:
        col_widths = [w / n_cols] * n_cols

    shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(l), Inches(t),
        Inches(w), Inches(total_h))
    tbl = shape.table

    # Disable built-in table styles so our fills show through
    tbl.first_row = False
    tbl.banded_rows = False

    # Column widths
    for j, cw in enumerate(col_widths):
        tbl.columns[j].width = Inches(cw)

    # Row heights
    tbl.rows[0].height = Inches(hdr_h)
    for i in range(1, n_rows):
        tbl.rows[i].height = Inches(row_h)

    def write_cell(cell, text, bg, fg, size, bold=False, align=PP_ALIGN.LEFT):
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        cell.margin_left   = Inches(0.06)
        cell.margin_right  = Inches(0.06)
        cell.margin_top    = Inches(0.05)
        cell.margin_bottom = Inches(0.05)
        tf = cell.text_frame
        tf.word_wrap = True
        # Clear any default empty run from the first paragraph
        first_p = tf.paragraphs[0]
        for run in first_p.runs:
            run._r.getparent().remove(run._r)
        lines = str(text).split('\n')
        for pi, line in enumerate(lines):
            p = first_p if pi == 0 else tf.add_paragraph()
            p.alignment = align
            r = p.add_run()
            r.text = line
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = fg

    # Header row
    for j, hdr in enumerate(headers):
        write_cell(tbl.rows[0].cells[j], hdr,
                   TABLE_HDR, TITLE_FG, hdr_sz,
                   bold=True, align=PP_ALIGN.CENTER)

    # Data rows
    for ri, row in enumerate(rows):
        bg = TABLE_ALT if ri % 2 == 0 else TABLE_REG
        for j, text in enumerate(row):
            bold = str(text).startswith("**") and str(text).endswith("**")
            txt = str(text).strip("*")
            fg = ACCENT if bold else BODY_FG
            write_cell(tbl.rows[ri + 1].cells[j], txt,
                       bg, fg, row_sz, bold=bold)

    _clear_cell_borders(tbl)
    return shape


def img(slide, path, l, t, w, max_h=None, caption=None):
    """Embed image at correct aspect ratio.

    w        — desired width in inches
    max_h    — optional maximum height; if the natural height exceeds this,
               the image is scaled down proportionally
    Returns the actual (width, height) placed.
    """
    p = Path(path)
    if p.exists():
        pil = PILImage.open(str(p))
        iw, ih = pil.size
        ratio = iw / ih
        actual_w = w
        actual_h = w / ratio
        if max_h and actual_h > max_h:
            actual_h = max_h
            actual_w = max_h * ratio
        slide.shapes.add_picture(
            str(p),
            Inches(l), Inches(t),
            Inches(actual_w), Inches(actual_h))
        if caption:
            txb(slide, caption,
                l, t + actual_h + 0.06,
                actual_w, 0.3,
                size=10, italic=True, color=DIM, align=PP_ALIGN.CENTER)
        return actual_w, actual_h
    else:
        bh = max_h or 2.0
        rect(slide, l, t, w, bh,
             fill=RGBColor(0x1A, 0x28, 0x3C), line_color=DIM)
        txb(slide, f"[{p.name}]",
            l + 0.1, t + bh / 2 - 0.2, w - 0.2, 0.4,
            size=11, color=DIM, align=PP_ALIGN.CENTER)
        if caption:
            txb(slide, caption, l, t + bh + 0.06, w, 0.3,
                size=10, italic=True, color=DIM, align=PP_ALIGN.CENTER)
        return w, bh


def divider(text, sub=""):
    """Section title slide."""
    sl = prs.slides.add_slide(BLANK)
    rect(sl, 0, 0, 13.33, 7.5, fill=BG)
    rect(sl, 0, 3.1, 13.33, 0.055, fill=RULE_CLR)
    txb(sl, text, 1.0, 2.1, 11.33, 1.3,
        size=46, bold=True, color=TITLE_FG, align=PP_ALIGN.CENTER)
    if sub:
        txb(sl, sub, 1.0, 3.4, 11.33, 0.8,
            size=20, color=BODY_FG, align=PP_ALIGN.CENTER, italic=True)


def new_slide(title, section=""):
    sl = prs.slides.add_slide(BLANK)
    chrome(sl, title, section)
    return sl


def rule(slide, y):
    rect(slide, CX, y, CW, 0.03, fill=DIM)


def arrow_note(slide, text, y=6.6):
    rule(slide, y - 0.12)
    txb(slide, f"→  {text}", CX, y, CW, 0.55,
        size=14, italic=True, color=ACCENT)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
rect(sl, 0, 0, 13.33, 7.5, fill=BG)
rect(sl, 0, 2.7, 13.33, 0.055, fill=RULE_CLR)

txb(sl, "Multi-Label Toxic Comment Classification",
    0.8, 0.75, 11.73, 1.5,
    size=42, bold=True, color=TITLE_FG, align=PP_ALIGN.CENTER)

txb(sl, "From CNN+GloVe to BERT\n"
        "Architecture · Data Efficiency · Rare Label Detection",
    0.8, 2.0, 11.73, 1.0,
    size=20, italic=True, color=BODY_FG, align=PP_ALIGN.CENTER)

txb(sl, "CMPE 258 — Deep Learning   |   Kaggle Jigsaw Toxic Comment Classification Challenge",
    0.8, 2.95, 11.73, 0.5,
    size=14, color=DIM, align=PP_ALIGN.CENTER)

txb(sl, '"Online platforms generate millions of comments daily.\n'
        'Human moderation doesn\'t scale — but can a model learn\n'
        'six different kinds of harm simultaneously?"',
    1.5, 3.6, 10.33, 1.5,
    size=18, italic=True, color=BODY_FG, align=PP_ALIGN.CENTER)

txb(sl, "[Name 1]   ·   [Name 2]   ·   [Name 3]   ·   [Name 4]",
    0.8, 5.7, 11.73, 0.55,
    size=17, color=BODY_FG, align=PP_ALIGN.CENTER)
txb(sl, "Spring 2026", 0.8, 6.25, 11.73, 0.4,
    size=13, color=DIM, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Research Question
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("Research Question")

txb(sl,
    "How much does model architecture and pretraining matter\n"
    "for multi-label toxicity detection — and how much training data\n"
    "do you actually need?",
    CX, 1.25, CW, 1.5,
    size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

rule(sl, 2.85)

txb(sl, "Three sub-questions that drove the full experiment chain:",
    CX, 3.0, CW, 0.45, size=16, color=DIM)

blist(sl, [
    (0, "Can task-trained word-vector models (CNN, BiLSTM) achieve acceptable performance on rare toxicity labels?"),
    (0, "Do pretrained transformers (DistilBERT, BERT) overcome the rare-label ceiling — and why?"),
    (0, "What is the minimum training data needed for each model family to reach competitive performance?"),
], CX, 3.5, CW, 3.0, size=18, gap_pt=14)

txb(sl, "6 labels:  toxic  ·  severe_toxic  ·  obscene  ·  threat  ·  insult  ·  identity_hate",
    CX, 6.7, CW, 0.45,
    size=15, color=DIM, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Dataset Overview
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("Dataset Overview", "DATA")

blist(sl, [
    (0, "159,571 Wikipedia talk-page comments  —  no missing values, no duplicate IDs"),
    (0, "Multi-label binary classification: each comment independently belongs to any subset of 6 toxicity categories"),
    (0, "Labels are NOT mutually exclusive — a single comment can be toxic, obscene, and an insult simultaneously"),
    (0, "Fixed train / val split held constant across all experiments:"),
    (1, "Train:       143,613 samples"),
    (1, "Validation:   15,958 samples   (iterative stratification, random_state=42)"),
    (0, "6 binary label columns:  toxic  ·  severe_toxic  ·  obscene  ·  threat  ·  insult  ·  identity_hate"),
], CX, CY, CW, 5.2, size=18, gap_pt=14)

rule(sl, 6.6)
txb(sl,
    "The same 15,958-sample validation set is used in every single experiment across all four models.",
    CX, 6.65, CW, 0.45, size=14, italic=True, color=DIM)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Label Distribution
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("Label Distribution", "EDA")

txb(sl, "~89.8% of comments are completely clean  ·  only ~10.2% have at least one label",
    CX, 1.25, CW, 0.45,
    size=18, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

# Table on left side
native_table(sl,
    ["Label", "Count", "Prevalence", "Note"],
    [
        ["toxic",         "~15,294", "~9.6%",   "Most common"],
        ["obscene",       "~8,457",  "~5.3%",   ""],
        ["insult",        "~7,819",  "~4.9%",   ""],
        ["severe_toxic",  "1,595",   "~1.0%",   "Always co-occurs with toxic"],
        ["identity_hate", "~1,405",  "~0.88%",  ""],
        ["**threat**",    "**478**", "**~0.30%**", "**Rarest — 32:1 ratio vs toxic**"],
    ],
    CX, 1.82, 7.8,
    col_widths=[2.1, 1.6, 1.6, 2.5],
    hdr_sz=15, row_sz=14, row_h=0.50, hdr_h=0.52)

# Image on right — fig01 is 1793×762 (ratio 2.35); w=5.2 → h≈2.21"
# Center vertically in content area (1.3 to 7.2): mid = 4.25, top = 4.25 - 1.1 = 3.15
img(sl, FIG / "fig01_class_distribution.png",
    8.5, 2.8, 4.6, max_h=2.8,
    caption="Per-label class prevalence")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Multi-label Co-occurrence
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("Label Co-occurrence", "EDA")

txb(sl, "Among the 16,225 labeled comments — how many toxicity labels does each carry?",
    CX, 1.25, CW, 0.45, size=17, color=DIM)

# Table on left
native_table(sl,
    ["# Labels Active", "Comment Count"],
    [
        ["1 label",  "6,360"],
        ["2 labels", "3,480"],
        ["3 labels", "4,209"],
        ["4 labels", "1,760"],
        ["5 labels", "385"],
        ["6 labels (all)", "31"],
    ],
    CX, 1.82, 4.6,
    col_widths=[2.3, 2.3],
    hdr_sz=15, row_sz=15, row_h=0.5, hdr_h=0.52)

# Key co-occurrence findings on right
txb(sl, "Key Co-occurrence Findings", 5.5, 1.82, 7.5, 0.5,
    size=16, bold=True, color=TITLE_FG)

blist(sl, [
    (0, "severe_toxic ⊂ toxic: 0 rows with severe_toxic=1 and toxic=0"),
    (0, "toxic – obscene – insult form a tight cluster:  Pearson ρ ≈ 0.65–0.74"),
    (0, "P(toxic=1 | obscene=1) = 93.81%"),
    (0, "P(toxic=1 | insult=1) = 93.23%"),
    (0, "P(toxic=1 | obscene=1 AND insult=1) = 96.73%"),
    (0, "threat weakly correlates with all others  (ρ ≈ 0.12–0.16)"),
    (1, "High average AUC can completely hide a full threat failure"),
], 5.5, 2.38, 7.5, 4.5, size=16, gap_pt=10)

rule(sl, 6.6)
txb(sl, "Source: EDA.ipynb — Pearson correlation heatmap + conditional probability heatmap P(B=1|A=1)",
    CX, 6.65, CW, 0.4, size=12, italic=True, color=DIM)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Text Characteristics
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("Text Length Characteristics", "EDA")

# Table on left
native_table(sl,
    ["Statistic", "Characters", "Words"],
    [
        ["Median",   "207", "36"],
        ["Mean",     "397", "67"],
        ["90th pct", "895", "152"],
        ["95th pct", "1,364", "230"],
        ["99th pct", "3,471", "567"],
        ["Max",      "5,000", "1,411"],
    ],
    CX, 1.3, 5.8,
    col_widths=[2.0, 1.9, 1.9],
    hdr_sz=15, row_sz=15, row_h=0.50, hdr_h=0.52)

# Findings on right side — full height
txb(sl, "Design Implications", 6.7, 1.3, 6.4, 0.5,
    size=17, bold=True, color=TITLE_FG)

blist(sl, [
    (0, "Toxic comments are shorter at the median (~94–123 chars) vs. clean (~206–216 chars)"),
    (1, "Length is a misleading proxy — do not use it as a standalone feature"),
    (0, "Heavy right tail: max 1,411 words — set max_length deliberately for all sequence models"),
    (0, "Linguistic noise: slang, misspellings, abbreviations"),
    (1, "Affects word-level tokenisers (CNN, BiLSTM) more than WordPiece (BERT, DistilBERT)"),
    (0, "Practical limits chosen:"),
    (1, "max_len = 100  (BiLSTM)  — covers 36-word median with headroom"),
    (1, "max_len = 256  (CNN)"),
    (1, "max_length = 128 → 192  (BERT/DistilBERT)  — key design decision"),
], 6.7, 1.9, 6.4, 5.0, size=16, gap_pt=10)

rule(sl, 6.6)
txb(sl, "Source: EDA.ipynb — comment length KDE (clean vs. any-toxic) + median char length by label value",
    CX, 6.65, CW, 0.4, size=12, italic=True, color=DIM)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — EDA → Design Decisions
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("EDA → Design Decisions", "EDA")

native_table(sl,
    ["EDA Finding", "Design Response"],
    [
        ["32:1 imbalance (toxic vs. threat)",
         "pos_weight per label in BCE loss  +  per-label threshold tuning"],
        ["threat has only 478 positives in 143k rows",
         "Training size sweep 10k→143k  ·  Macro F1 as primary metric"],
        ["severe_toxic ⊂ toxic on train",
         "Separate head + own threshold; expect unstable F1"],
        ["threat weakly correlated with other labels",
         "Evaluate threat explicitly — average AUC hides failure"],
        ["Length variance: 36 median → 1,411 max words",
         "max_len=100 (BiLSTM)  ·  256 (CNN)  ·  128→192 (BERT)"],
        ["66% of this vocabulary has no GloVe vector",
         "WordPiece subword tokenisation for transformers handles OOV via decomposition"],
        ["Label co-occurrence structure in splits",
         "Iterative stratification (skmultilearn) instead of random split"],
    ],
    CX, 1.3, CW,
    col_widths=[5.0, 7.33],
    hdr_sz=15, row_sz=14, row_h=0.54, hdr_h=0.52)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Methodology: Metrics
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("Overall Methodology — Metrics", "METHODOLOGY")

native_table(sl,
    ["Metric", "What It Is", "Why We Use It"],
    [
        ["Tuned Macro F1\n(PRIMARY)",
         "Unweighted average of per-label F1\nacross all 6 labels, after threshold optimisation",
         "Weights all 6 labels equally — a complete failure on threat penalises the score the same as a failure on toxic"],
        ["Baseline Macro F1",
         "Same metric but with fixed threshold\nt=0.5 for all labels",
         "Measures raw model calibration before any threshold adjustment"],
        ["Tuning Gain\n(Tuned − Baseline)",
         "Difference between tuned and\nbaseline macro F1",
         "Large gain = model discriminates well but outputs miscalibrated probabilities"],
        ["ROC-AUC per label\n(secondary)",
         "Probability the model ranks a positive\nabove a negative — threshold-independent",
         "Confirms whether poor F1 is a calibration problem or a discrimination problem"],
        ["Per-Label F1",
         "F1 for each of the 6 labels individually",
         "Exposes which specific toxicity types each model struggles with"],
    ],
    CX, 1.3, CW,
    col_widths=[2.4, 3.8, 6.13],
    hdr_sz=14, row_sz=13, row_h=0.66, hdr_h=0.5)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Methodology: Types of Runs
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("Overall Methodology — Types of Runs", "METHODOLOGY")

native_table(sl,
    ["Run Type", "What It Does", "Which Models"],
    [
        ["Full-data training run",
         "Train on all 143,613 samples; evaluate on\nfixed 15,958-sample validation set",
         "All 4 models"],
        ["Training size sweep",
         "Train at 10k, 20k, 30k … up to full data;\nsame fixed val set at every size",
         "CNN (→61k)  ·  BiLSTM (→140k)\nDistilBERT & BERT (→143k)"],
        ["Hyperparameter grid search",
         "Hold training size fixed at full data;\nvary lr / weight decay / warmup across 16 configs",
         "BERT Exp04 only"],
        ["Ablation study",
         "Hold all settings fixed except one variable;\nisolate and measure the effect of that single change",
         "DistilBERT Exp D (pos_weight)\nDistilBERT Exp E (loss function)"],
        ["Per-label threshold tuning",
         "Grid search over [0.05–0.995] per label\non the validation set after every single run",
         "Applied to ALL runs, ALL models"],
    ],
    CX, 1.3, CW,
    col_widths=[2.9, 5.5, 3.93],
    hdr_sz=14, row_sz=13, row_h=0.68, hdr_h=0.5)

rule(sl, 6.75)
txb(sl, "Fixed anchor: the same 15,958-sample validation set is used in every run — this is what makes all sweep curves directly comparable.",
    CX, 6.8, CW, 0.45, size=13, italic=True, color=DIM)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — CNN Architecture
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("CNN + GloVe — Architecture", "CNN + GLOVE")

# Text on left — 6 key bullets
blist(sl, [
    (0, "TextCNN: parallel Conv1d filters for 2 / 3 / 4 / 5-gram patterns"),
    (1, "128 filters per filter size  ·  global max-pool  ·  512-dim concatenated output  ·  Dropout(0.5)"),
    (0, "Parameters scale with vocabulary size:"),
    (1, "~2.6M at 10k training samples  (vocab = 24,700 unique tokens)"),
    (1, "~5.1M at max_vocab = 50k  (full training run)"),
    (0, "Embeddings: GloVe 6B 100-dim — fine-tuned during training  (GLOVE_TRAINABLE = True)"),
    (1, "Only 34% vocabulary coverage — 66% initialise from random vectors"),
    (0, "max_len=256  ·  max_vocab=50k  ·  min_freq=2  ·  BCEWithLogitsLoss + pos_weight"),
], CX, CY, CW, 3.5, size=17, gap_pt=10)

# graph8_vocab_coverage is 2086×668 (ratio 3.12)
# At w=10.0 → h=10.0/3.12=3.21". Center x=(13.33-10)/2=1.67
# Place below text at y=5.0, ends at 8.21 — too tall.
# Use w=9.0 → h=2.88", x=2.17, y=4.85, ends at 7.73 — still tight.
# Use w=8.5 → h=2.72", x=2.42, y=4.85, ends at 7.57 — close.
# Use max_h=2.5 with w=9.0 → scales to w=9.0,h=2.5 if h>max_h: h=2.88>2.5, so w=2.5*3.12=7.8, h=2.5
# Place at y=4.8, x=(13.33-7.8)/2=2.77, ends at 7.3.
img(sl, NB / "cnn_glove/cnn_glove_sweep/graph8_vocab_coverage.png",
    2.77, 4.8, 9.0, max_h=2.4,
    caption="GloVe vocabulary coverage vs. training size  (graph8_vocab_coverage.png)")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CNN Sweep & Outcomes
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("CNN + GloVe — Sweep & Outcomes", "CNN + GLOVE")

# Table on left
txb(sl, "Full-data result (143,613 samples, rebalance_train=False):",
    CX, 1.3, 6.2, 0.4, size=14, color=DIM)

native_table(sl,
    ["Label", "Tuned F1", "ROC-AUC"],
    [
        ["toxic",         "0.753", "0.961"],
        ["severe_toxic",  "0.498", "0.985"],
        ["obscene",       "0.772", "0.980"],
        ["**threat**",    "**0.277**", "0.982"],
        ["insult",        "0.687", "0.971"],
        ["identity_hate", "0.379", "0.963"],
        ["**Macro F1**",  "**0.561**", "—"],
    ],
    CX, 1.76, 5.0,
    col_widths=[2.3, 1.35, 1.35],
    hdr_sz=14, row_sz=14, row_h=0.46, hdr_h=0.48)

blist(sl, [
    (0, "Training time: 886s on CPU"),
    (0, "ROC-AUC is 0.96+ on all labels — model ranks toxic vs. clean well"),
    (0, "F1 for threat (0.277) and identity_hate (0.379) is very poor"),
    (1, "66% OOV: slurs and toxic slang start from random init with very few examples"),
    (0, "Bottleneck is representation quality, not architecture"),
], CX, 5.35, 5.8, 2.0, size=16, gap_pt=8)

# graph1_f1_vs_size is 1782×730 (ratio 2.44)
# Right side: x=6.3, w=6.7 → h=6.7/2.44=2.74"
# Center vertically: content area 1.3→7.2 (5.9"), image h=2.74", top = 1.3+(5.9-2.74)/2 = 2.88
img(sl, NB / "cnn_glove/cnn_glove_sweep/graph1_f1_vs_size.png",
    6.5, 2.5, 6.6, max_h=3.0,
    caption="Macro F1 vs. training size  (Colab sweep, rebalance_train=True)")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 12 — BiLSTM Architecture
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("BiLSTM + Attention — Architecture", "BiLSTM")

blist(sl, [
    (0, "Packed sequences — handles variable-length comments without padding overhead"),
    (0, "Spatial Dropout(0.2) on the embedding dimension — regularises at representation level"),
    (0, "Bidirectional LSTM (hidden=128) — reads sequence left-to-right and right-to-left simultaneously"),
    (0, "LayerNorm after BiLSTM  ·  Learned attention with padding mask"),
    (1, "Attention weights each token's contribution — not limited to fixed n-gram windows like CNN"),
    (0, "~2M parameters  ·  max_len=100  (EDA: covers the ~36-word median with ample headroom)"),
    (0, "Learnable embeddings — NOT GloVe — the model learns toxic vocabulary directly from context"),
    (1, "Completely eliminates the 66% OOV problem — there are no unknown tokens"),
    (0, "BCEWithLogitsLoss + pos_weight  ·  early stopping on val macro F1  (patience=3)"),
    (0, "Training size sweep: 20k → 140k  (no 10k data point — sweep started at 20k)"),
], CX, CY, CW, 5.5, size=17, gap_pt=10)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 13 — BiLSTM Outcomes
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("BiLSTM + Attention — Outcomes", "BiLSTM")

txb(sl, "Best result at 60k training samples  (performance degraded beyond 60k — overfitting to noisy labels):",
    CX, 1.3, CW, 0.4, size=14, color=DIM)

native_table(sl,
    ["Label", "Tuned F1", "ROC-AUC"],
    [
        ["toxic",         "0.737", "0.961"],
        ["severe_toxic",  "0.457", "0.985"],
        ["obscene",       "0.745", "0.968"],
        ["**threat**",    "**0.462**", "0.998"],
        ["insult",        "0.659", "0.960"],
        ["identity_hate", "0.270", "0.953"],
        ["**Macro F1**",  "**0.555**", "—"],
    ],
    CX, 1.78, 5.0,
    col_widths=[2.3, 1.35, 1.35],
    hdr_sz=14, row_sz=14, row_h=0.50, hdr_h=0.48)

blist(sl, [
    (0, "Threat F1 (0.462) much better than CNN (0.277) — learnable embeddings specialise on toxic vocab"),
    (0, "Identity_hate regressed vs. CNN (0.270 vs. 0.379) — fewer positives at 60k than full 143k"),
    (0, "Macro F1 (0.555) ≈ CNN (0.561) despite very different architectures"),
    (0, "Key finding: both models hit the same ceiling — word-level representation is the bottleneck"),
    (0, "Motivates transformers: WordPiece subword tokenisation handles OOV via decomposition — no true unknowns"),
], CX, 5.7, CW, 2.0, size=16, gap_pt=10)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 14 — DistilBERT Exp B
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("DistilBERT — Exp B: Baseline", "DISTILBERT")

blist(sl, [
    (0, "distilbert-base-uncased: 6 transformer layers  ·  768 hidden  ·  12 heads  ·  ~67M params"),
    (0, "WordPiece subword tokenisation: decomposes unknown words into known subwords — eliminates OOV entirely"),
    (0, "Config: max_length=128  ·  lr=2e-5  ·  warmup=0.1  ·  wd=0.01  ·  pos_weight=True  ·  BCE"),
], CX, 1.3, CW, 1.6, size=17, gap_pt=10)

txb(sl, "Training size sweep results:", CX, 3.05, 5.0, 0.4, size=15, color=DIM)

native_table(sl,
    ["Train Size", "Tuned Macro F1"],
    [
        ["10k",  "0.584"],
        ["40k",  "0.659"],
        ["60k",  "0.671"],
        ["143k", "**0.676**"],
    ],
    CX, 3.52, 5.0,
    col_widths=[2.5, 2.5],
    hdr_sz=14, row_sz=15, row_h=0.54, hdr_h=0.5)

blist(sl, [
    (0, "Full data: macro F1 = 0.676  vs. CNN 0.561  →  +0.115 purely from representation quality"),
    (0, "Curve is noisy — drops at 20k→30k before recovering — pos_weight may be fighting pretrained priors"),
], CX, 5.8, CW, 1.0, size=16, gap_pt=10)

arrow_note(sl, "Next: ablate pos_weight to isolate its effect on a pretrained model")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 15 — DistilBERT Exp D
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("DistilBERT — Exp D: Positive Weighting Ablation", "DISTILBERT")

txb(sl, "Question: does pos_weight help or hurt DistilBERT?",
    CX, 1.3, CW, 0.5, size=20, bold=True, color=ACCENT)

txb(sl, "6-variant ablation at full data — pos_weight × {none, clip20, clip50} × {baseline, rebalanced}:",
    CX, 1.88, CW, 0.42, size=14, color=DIM)

native_table(sl,
    ["Config", "Tuned Macro F1"],
    [
        ["**no pos_weight (baseline)**", "**0.7007**"],
        ["rebalanced + pw_clip20",       "0.6958"],
        ["rebalanced + no_pw",           "0.6923"],
        ["pos_weight clip20",            "0.6923"],
        ["pos_weight clip50",            "0.6906"],
        ["rebalanced + pw_clip50",       "0.6887"],
    ],
    CX, 2.38, 7.0,
    col_widths=[4.5, 2.5],
    hdr_sz=15, row_sz=15, row_h=0.54, hdr_h=0.5)

blist(sl, [
    (0, "Finding: pos_weight hurts DistilBERT — +0.024 gain from simply removing it"),
    (0, "Pretrained models have internal class priors from pretraining on massive corpora"),
    (1, "Forcing external pos_weight fights those priors and distorts calibration"),
    (0, "Decision: drop pos_weight for all subsequent transformer experiments"),
], CX, 5.85, CW, 1.5, size=17, gap_pt=10)

arrow_note(sl, "Next: test whether focal loss can further improve rare-label calibration")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 16 — DistilBERT Exp E
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("DistilBERT — Exp E: Loss Function Comparison", "DISTILBERT")

txb(sl, "Question: can focal loss beat BCE for class-imbalanced multi-label classification?",
    CX, 1.3, CW, 0.5, size=20, bold=True, color=ACCENT)

txb(sl, "Config: lr=1.75e-5  ·  wd=0.015  ·  warmup=0.1  ·  no pos_weight  ·  full data (143k):",
    CX, 1.88, CW, 0.42, size=14, color=DIM)

native_table(sl,
    ["Loss Function", "Tuned Macro F1"],
    [
        ["**BCE, no pos_weight**", "**0.7007**"],
        ["Focal  γ=3",             "0.6965"],
        ["Focal  γ=1",             "0.6923"],
        ["Focal  γ=2",             "0.6870"],
    ],
    CX, 2.38, 6.5,
    col_widths=[4.0, 2.5],
    hdr_sz=15, row_sz=15, row_h=0.56, hdr_h=0.5)

blist(sl, [
    (0, "Finding: focal loss underperforms BCE across all γ values"),
    (0, "Focal loss and per-label threshold tuning both address calibration — using both is redundant"),
    (1, "Focal loss adjusts class weights during training"),
    (1, "Threshold tuning adjusts the decision boundary after training"),
    (1, "Stacking them attacks the same problem from both ends — they conflict"),
    (0, "Decision: BCE + per-label threshold tuning is the final recipe for all transformer experiments"),
], CX, 4.72, CW, 2.3, size=17, gap_pt=10)

arrow_note(sl, "Apply best config to full training-size sweep to measure DistilBERT data efficiency")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 17 — DistilBERT Final Sweep
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("DistilBERT — Final Sweep & Outcomes", "DISTILBERT")

txb(sl, "Best config: no pos_weight  ·  BCE  ·  lr=1.75e-5  ·  wd=0.015  ·  warmup=0.1  ·  max_length=128",
    CX, 1.3, 7.0, 0.42, size=14, color=DIM)

native_table(sl,
    ["Train Size", "Tuned Macro F1"],
    [
        ["10k",  "0.544"],
        ["40k",  "0.660"],
        ["50k",  "0.679"],
        ["80k",  "0.689"],
        ["90k",  "**0.698**"],
        ["140k", "0.704"],
        ["143k", "0.696"],
    ],
    CX, 1.8, 5.0,
    col_widths=[2.5, 2.5],
    hdr_sz=14, row_sz=14, row_h=0.5, hdr_h=0.5)

blist(sl, [
    (0, "Full data (143k): tuned macro F1 = 0.696  ·  training time = 1,403s on GPU"),
    (0, "+0.020 over Exp B (0.676) purely from dropping pos_weight"),
    (0, "Curve peaks at 140k (0.704) — noisy above 80k with no clean plateau"),
    (0, "Rapid gains 10k→50k; strong diminishing returns beyond 80k"),
], CX, 5.55, 7.0, 1.8, size=16, gap_pt=10)

# fig07 is 1934×808 (ratio 2.39); w=5.3 → h=2.22"; center y=(1.3+7.2)/2 - 1.11 = 4.25-1.11=3.14
img(sl, FIG / "fig07_distilbert_per_label_vs_size.png",
    7.7, 2.6, 5.3,
    caption="DistilBERT per-label F1 vs. training size")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 18 — BERT Exp04
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("BERT — Exp04: Hyperparameter Grid Search", "BERT")

blist(sl, [
    (0, "bert-base-uncased: 12 transformer layers  ·  768 hidden  ·  12 heads  ·  ~109M params  ·  max_length=128"),
    (0, "16-run grid: lr ∈ {1.5e-5, 2.0e-5, 2.5e-5, 3.0e-5}  ×  wd ∈ {0.01, 0.015}  ×  warmup ∈ {0.06, 0.10}"),
    (0, "All runs: BCE, no pos_weight, full data (143k)"),
], CX, 1.3, 7.5, 1.5, size=17, gap_pt=10)

txb(sl, "Result range: 0.6986 – 0.7091   (spread of only 0.010)",
    CX, 2.95, 7.5, 0.48, size=19, bold=True, color=BODY_FG)
txb(sl, "Narrow spread — model quality drives performance more than hyperparameter tuning",
    CX, 3.42, 7.5, 0.42, size=14, italic=True, color=DIM)

txb(sl, "Best config:  lr=1.5e-5  ·  wd=0.015  ·  warmup=0.06   →   Macro F1 = 0.7091",
    CX, 3.95, 7.5, 0.48, size=18, bold=True, color=ACCENT)

blist(sl, [
    (0, "Threat F1 = 0.537 — identified as the key bottleneck"),
    (0, "Hypothesis: max_length=128 may be truncating context in threat and identity_hate comments"),
    (1, "These comments often contain preamble before the actual threat — truncation cuts relevant signal"),
], CX, 4.55, 7.5, 1.8, size=17, gap_pt=10)

# fig06 is 1271×583 (ratio 2.18); w=5.0 → h=2.29"
# center y=(1.3+7.2)/2 - 2.29/2 = 4.25-1.15 = 3.10
img(sl, FIG / "fig06_bert_sweep_heatmap.png",
    8.1, 2.8, 5.0,
    caption="BERT Exp04 hyperparameter grid heatmap")

arrow_note(sl, "Next: increase context window 128 → 192 tokens and measure the per-label impact")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 19 — BERT Exp05
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("BERT — Exp05: Context Window Hypothesis", "BERT")

txb(sl, "Hypothesis: increasing max_length 128 → 192 captures more context for threat and identity_hate",
    CX, 1.3, CW, 0.5, size=19, bold=True, color=ACCENT)
txb(sl, "Best hyperparams from Exp04  ·  bf16 AMP  ·  grad_accum_steps=2, batch=16  (effective batch=32)",
    CX, 1.88, CW, 0.4, size=13, color=DIM)

native_table(sl,
    ["Label", "Exp04  (128 tok)", "Exp05  (192 tok)", "Δ"],
    [
        ["**threat**",        "0.537", "**0.659**", "**+0.122**"],
        ["**identity_hate**", "0.534", "**0.597**", "**+0.063**"],
        ["severe_toxic",      "0.533", "0.552",     "+0.019"],
        ["toxic",             "0.839", "0.841",     "+0.002"],
        ["obscene",           "—",     "0.849",     "—"],
        ["insult",            "—",     "0.784",     "—"],
        ["**Macro F1**",      "0.709", "**0.714**", "+0.005"],
    ],
    CX, 2.36, 8.5,
    col_widths=[2.3, 2.0, 2.0, 2.2],
    hdr_sz=14, row_sz=14, row_h=0.5, hdr_h=0.5)

blist(sl, [
    (0, "Full data: tuned macro F1 = 0.7135  —  best single BERT result"),
    (0, "Training time = 2,937s on GPU (L4)"),
    (0, "Context window is a first-class design parameter for rare-label classification"),
], CX, 6.0, 8.5, 1.0, size=16, gap_pt=10)

arrow_note(sl, "Apply Exp05 config to full training-size sweep (Exp06) to measure BERT data efficiency")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 20 — BERT Exp06
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("BERT — Exp06: Data Efficiency Sweep", "BERT")

txb(sl, "Same config as Exp05  (max_length=192, lr=1.5e-5, wd=0.015, warmup=0.06)  ·  15 training sizes:",
    CX, 1.3, 7.5, 0.42, size=14, color=DIM)

native_table(sl,
    ["Train Size", "Tuned Macro F1", "Train Time"],
    [
        ["10k",   "0.591",      "254s"],
        ["50k",   "0.690",      "1,274s"],
        ["80k",   "0.696",      "2,041s"],
        ["90k",   "**0.701**",  "1,837s"],
        ["120k",  "**0.710**",  "2,450s"],
        ["143k",  "0.706",      "2,944s"],
    ],
    CX, 1.8, 6.5,
    col_widths=[2.1, 2.2, 2.2],
    hdr_sz=14, row_sz=15, row_h=0.54, hdr_h=0.5)

blist(sl, [
    (0, "BERT at 80k (0.696) = DistilBERT at full data (0.696) — same macro F1 at 56% of the data"),
    (0, "Peak at 120k (0.710) — full data (0.706) is not the single best point"),
    (0, "50k → full data gain = only +0.016 — strong diminishing returns above 50k"),
    (0, "~91% of total macro F1 gain (10k → full data) is captured at just 80k samples"),
], CX, 5.45, 7.5, 2.0, size=16, gap_pt=10)

# bert_exp_06_graph_01 is 1650×750 (ratio 2.20); w=5.1 → h=2.32"
# center y=(1.3+7.2)/2 - 2.32/2 = 4.25-1.16 = 3.09
img(sl, NB / "bert/bert_06/bert_exp_06_graph_01_f1_vs_size.png",
    8.1, 2.7, 5.0,
    caption="BERT Exp06 — Macro F1 vs. training size")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Per-Label F1 Comparison
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("Model Comparison — Per-Label F1 at Full Data", "RESULTS")

native_table(sl,
    ["Label", "CNN", "BiLSTM\n(best @ 60k)", "DistilBERT", "BERT Exp05"],
    [
        ["toxic",         "0.753", "0.737", "0.841", "0.841"],
        ["severe_toxic",  "0.498", "0.457", "0.543", "0.552"],
        ["obscene",       "0.772", "0.745", "0.845", "0.849"],
        ["**threat**",    "0.277", "0.462", "0.568", "**0.659**"],
        ["identity_hate", "0.379", "0.270", "**0.608**", "0.597"],
        ["insult",        "0.687", "0.659", "0.772", "0.784"],
        ["**Macro F1**",  "0.561", "0.555", "0.696", "**0.714**"],
    ],
    CX, 1.3, CW,
    col_widths=[2.3, 1.9, 2.1, 2.0, 2.03],
    hdr_sz=14, row_sz=14, row_h=0.54, hdr_h=0.52)

blist(sl, [
    (0, "BERT Exp05 leads overall (0.714) and on threat (0.659) — context window extension was decisive"),
    (0, "DistilBERT beats BERT on identity_hate (0.608 vs. 0.597) — smaller model, same context window"),
    (0, "Both task-trained models plateau at ~0.56 despite different architectures — representation ceiling"),
], CX, 6.1, CW, 1.2, size=16, gap_pt=10)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Comparison at 10k
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("Model Comparison — At 10k Training Samples", "RESULTS")

txb(sl,
    "How well does each model perform when given only 10,000 training examples?",
    CX, 1.3, CW, 0.5, size=20, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)

native_table(sl,
    ["Model", "Macro F1 @ 10k", "Note"],
    [
        ["**BERT Exp06**",        "**0.591**", "Best at low data — pretraining transfers immediately"],
        ["DistilBERT (final)",    "0.544",     "Strong from 10k onward"],
        ["CNN (rebalanced sweep)", "0.533",    "Rebalanced preprocessing — not directly comparable"],
        ["BiLSTM",                "no 10k data", "Sweep started at 20k"],
    ],
    CX, 2.0, CW,
    col_widths=[3.2, 2.2, 6.93],
    hdr_sz=15, row_sz=15, row_h=0.58, hdr_h=0.52)

blist(sl, [
    (0, "BERT leads at all observed training sizes — there is no crossover point where task-trained models win"),
    (0, "Both transformers already outperform CNN at 10k samples — pretraining transfers immediately"),
    (0, "CNN's 10k value uses rebalanced data (rebalance_train=True); BERT/DistilBERT use raw distribution"),
    (1, "Direct numeric comparison is inexact — CNN had access to fewer negative examples"),
    (0, "Practical implication: if you can afford even minimal GPU fine-tuning, a pretrained transformer always wins"),
], CX, 4.5, CW, 2.8, size=17, gap_pt=10)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Data Efficiency Curves
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("Data Efficiency — All Models", "RESULTS")

blist(sl, [
    (0, "CNN and BiLSTM plateau early — CNN peaks ~60k (macro F1 ≈ 0.58), no benefit from more data"),
    (0, "Both transformers show rapid gain 10k→50k; curve flattens above 80k"),
    (0, "BERT at 80k (0.696) = DistilBERT at full data (0.696) — same F1 at 56% of the data"),
    (0, "~91% of BERT's total macro F1 gain (10k → full data) is captured at 80k samples"),
], CX, 1.3, CW, 2.5, size=18, gap_pt=12)

# fig03 is 1934×808 (ratio 2.39)
# w=9.5 → h=9.5/2.39=3.97"; center x=(13.33-9.5)/2=1.92; at y=4.0, ends at 7.97 — too tall
# w=8.5 → h=3.56"; center x=2.42; at y=4.0, ends at 7.56 — just over
# w=8.0 → h=3.35"; center x=2.67; at y=4.0, ends at 7.35 — fits
# Use max_h=3.1 with w=9.0: h=9.0/2.39=3.77>3.1, so w=3.1*2.39=7.41, h=3.1; x=(13.33-7.41)/2=2.96
img(sl, FIG / "fig03_macro_f1_vs_size.png",
    2.67, 3.95, 8.0,
    caption="All-model tuned macro F1 vs. training size  (fig03)")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Efficiency Frontier
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("Training Cost vs. Performance", "RESULTS")

native_table(sl,
    ["Model", "Train Time @ Full Data", "Hardware", "Macro F1"],
    [
        ["CNN",           "886s",       "CPU",      "0.561"],
        ["DistilBERT",    "1,403s",     "GPU (L4)", "0.696"],
        ["**BERT Exp05**", "**2,937s**", "**GPU (L4)**", "**0.714**"],
    ],
    CX, 1.3, 9.5,
    col_widths=[2.5, 2.5, 2.2, 2.3],
    hdr_sz=15, row_sz=15, row_h=0.58, hdr_h=0.52)

txb(sl, "ROC-AUC note: all models score 0.95+ on all labels — the gap is entirely in calibration, not ranking ability.",
    CX, 3.2, 9.5, 0.5, size=14, italic=True, color=DIM)

# fig05 is 1784×957 (ratio 1.86)
# w=7.5 → h=7.5/1.86=4.03"; center x=(13.33-7.5)/2=2.92; at y=3.85, ends at 7.88 — over
# w=6.5 → h=6.5/1.86=3.49"; center x=3.42; at y=3.85, ends at 7.34 — fits
# Use max_h=3.2 with w=7.5: h=7.5/1.86=4.03>3.2, so w=3.2*1.86=5.95, h=3.2; x=(13.33-5.95)/2=3.69
img(sl, FIG / "fig05_efficiency_frontier.png",
    3.0, 3.85, 7.5, max_h=3.3,
    caption="Efficiency frontier: tuned macro F1 vs. training time  (fig05)")


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 25 — Key Findings
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("Key Findings", "CONCLUSION")

blist(sl, [
    (0, "Representation is the bottleneck for task-trained models"),
    (1, "CNN and BiLSTM both plateau at ~0.56 macro F1; 66% OOV is a hard ceiling no architecture change can overcome"),
    (0, "WordPiece subword tokenisation accounts for most of the transformer advantage"),
    (1, "DistilBERT's +0.115 gain over CNN comes from handling OOV via decomposition, not from model depth"),
    (0, "Context window is a first-class design parameter for rare-label classification"),
    (1, "BERT 128→192 tokens gave +0.122 threat F1 — the single largest gain from a single change in this project"),
    (0, "pos_weight and focal loss hurt pretrained transformers"),
    (1, "Pretrained priors interact with external loss reweighting — per-label threshold tuning alone is the right approach"),
    (0, "BERT leads at all training sizes including 10k — no regime where task-trained models win"),
    (0, "Diminishing returns above 50k — 91% of BERT's macro F1 gain is captured at 80k samples"),
], CX, CY, CW, 5.5, size=17, gap_pt=9)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 26 — Limitations & Future Work
# ═══════════════════════════════════════════════════════════════════════════
sl = new_slide("Limitations & Future Work", "CONCLUSION")

txb(sl, "Limitations", CX, 1.3, 6.5, 0.42,
    size=20, bold=True, color=TITLE_FG)
blist(sl, [
    (0, "CNN sweep used rebalance_train=True — cross-model data efficiency comparison is inexact"),
    (0, "No cross-dataset evaluation — Wikipedia talk-page → Twitter / Reddit transfer is unknown"),
    (0, "Per-label threshold tuning on the fixed val set may slightly overfit to that 15,958-sample distribution"),
    (0, "BiLSTM early stopping on val macro F1 vs. CNN fixed 2 epochs makes their comparison slightly unfair"),
], CX, 1.82, CW, 2.2, size=17, gap_pt=10)

rule(sl, 4.3)

txb(sl, "Future Work", CX, 4.4, 6.5, 0.42,
    size=20, bold=True, color=TITLE_FG)
blist(sl, [
    (0, "RoBERTa or DeBERTa — more robust pretraining for further rare-label gains"),
    (0, "Ensemble CNN + BERT as a fast-filter / deep-reranker pipeline"),
    (0, "Calibration via Platt scaling or temperature scaling — cleaner than brute-force threshold grid search"),
    (0, "Cross-lingual extension using mBERT or XLM-R for non-English toxic content"),
], CX, 4.9, CW, 2.2, size=17, gap_pt=10)


# ═══════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════
prs.save(str(OUT))
print(f"Saved  →  {OUT}")
print(f"Slides :  {len(prs.slides)}")
